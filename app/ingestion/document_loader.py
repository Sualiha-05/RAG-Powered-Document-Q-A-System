from pypdf import PdfReader
import docx
import openpyxl
import pandas as pd
from pptx import Presentation
from odf import text as odf_text, teletype
from odf.opendocument import load as odf_load

# ── Tier 1: plain-text-based formats ─────────────────────────────────────
PLAIN_TEXT_EXTENSIONS = {"txt", "md", "py", "java", "cpp", "c", "html", "css"}


def load_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text


def load_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def load_plain_text(file_path: str) -> str:
    """Handles .txt, .md, .py, .java, .cpp, .c, .html, .css — all read as raw text."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def load_odt(file_path: str) -> str:
    doc = odf_load(file_path)
    paragraphs = doc.getElementsByType(odf_text.P)
    return "\n".join(teletype.extractText(p) for p in paragraphs)


def load_xlsx(file_path: str) -> str:
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    text_parts = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text_parts.append(f"--- Sheet: {sheet_name} ---")
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) for cell in row if cell is not None]
            if row_values:
                text_parts.append(" | ".join(row_values))
    return "\n".join(text_parts)


def load_csv(file_path: str) -> str:
    df = pd.read_csv(file_path)
    text_parts = [f"Columns: {', '.join(df.columns.astype(str))}"]
    for _, row in df.iterrows():
        row_text = " | ".join(f"{col}: {val}" for col, val in row.items())
        text_parts.append(row_text)
    return "\n".join(text_parts)


def load_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_parts = []
    for i, slide in enumerate(prs.slides, start=1):
        text_parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        text_parts.append(line)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                text_parts.append(f"[Notes] {notes}")
    return "\n".join(text_parts)


LOADERS = {
    "pdf": load_pdf,
    "docx": load_docx,
    "odt": load_odt,
    "xlsx": load_xlsx,
    "csv": load_csv,
    "pptx": load_pptx,
}


def load_document(file_path: str) -> str:
    ext = file_path.lower().split(".")[-1]

    if ext in PLAIN_TEXT_EXTENSIONS:
        text = load_plain_text(file_path)
    elif ext in LOADERS:
        text = LOADERS[ext](file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    if not text or not text.strip():
        raise ValueError(
            "No extractable text found. This usually means the file is a scanned "
            "image or otherwise has no text layer (OCR isn't supported yet)."
        )
    return text